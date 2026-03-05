# backend/app/ingestion/ms_graph_ingest.py
import html
import re
import io
import requests
from datetime import datetime, timezone
from pptx import Presentation

from fastapi import APIRouter, Header, HTTPException
from pydantic import BaseModel
from typing import Optional, List

from app.db import SessionLocal
from app.models import ScrapedMessage

MAX_SUBJECT_LEN = 300
MAX_MESSAGE_TEXT_LEN = 60000

router = APIRouter()


# =========================
# Token helpers (Swagger-friendly)
# =========================
def _get_bearer_token(authorization: Optional[str]) -> Optional[str]:
    if not authorization:
        return None
    if not authorization.lower().startswith("bearer "):
        return None
    return authorization.split(" ", 1)[1].strip()


def _normalize_token(token: str) -> str:
    # remove spaces/newlines that break tokens when copying from UI
    return re.sub(r"\s+", "", (token or ""))

def _validate_ms_access_token(access_token: str) -> str:
    token = _normalize_token(access_token)
    if not token:
        raise Exception("Missing Microsoft access_token")

    # Validate by calling Graph instead of assuming JWT format
    s = requests.Session()
    s.trust_env = False
    r = s.get(
        "https://graph.microsoft.com/v1.0/me",
        headers={"Authorization": f"Bearer {token}"},
        timeout=20,
    )
    if not r.ok:
        raise Exception(f"Graph token check failed ({r.status_code}): {r.text}")

    return token


def _resolve_token(authorization: Optional[str], body_access_token: Optional[str]) -> str:
    header_token = _get_bearer_token(authorization)
    body_token = body_access_token

    # try header token first
    if header_token:
        try:
            return _validate_ms_access_token(header_token)
        except Exception:
            pass

    # fallback to body token
    return _validate_ms_access_token(body_token)


# =========================
# Text helpers
# =========================
def _strip_html(raw_html: str) -> str:
    if not raw_html:
        return ""
    raw_html = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", raw_html)
    text = re.sub(r"(?s)<[^>]*>", " ", raw_html)
    text = html.unescape(text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _parse_dt(dt_str: str):
    if not dt_str:
        return datetime.now(timezone.utc)
    try:
        return datetime.fromisoformat(dt_str.replace("Z", "+00:00"))
    except Exception:
        return datetime.now(timezone.utc)


def _truncate(text: str, max_len: int) -> str:
    if not text:
        return ""
    if len(text) <= max_len:
        return text
    return text[:max_len] + " ...[truncated]"


def _extract_text_from_pptx_bytes(file_bytes: bytes) -> str:
    try:
        f = io.BytesIO(file_bytes)
        prs = Presentation(f)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        return f"[Error parsing PPTX: {e}]"


# =========================
# Ingestion implementations
# =========================
def ingest_outlook_mail(
    access_token: str,
    owner_role="owner",
    max_results=50,
    allowed_senders: Optional[List[str]] = None,
):
    # Optional whitelist cleanup (Swagger sometimes sends "string")
    cleaned_allowed = []
    for s in (allowed_senders or []):
        if not s:
            continue
        s = s.strip()
        if not s or s.lower() == "string":
            continue
        cleaned_allowed.append(s.lower())
    allowed_senders_set = set(cleaned_allowed)

    # Clamp for safety
    try:
        max_results = int(max_results)
    except Exception:
        max_results = 50
    max_results = max(1, min(max_results, 200))

    url = "https://graph.microsoft.com/v1.0/me/messages"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Accept": "application/json",
        "ConsistencyLevel": "eventual",
    }
    params = {
        "$top": max_results,
        "$select": "id,subject,body,from,receivedDateTime,conversationId",
        "$orderby": "receivedDateTime DESC",
    }

    session = requests.Session()
    session.trust_env = False

    try:
        r = session.get(url, headers=headers, params=params, timeout=30)
        if not r.ok:
            raise Exception(f"Graph {r.status_code}: {r.text}")

        data = r.json()
        messages = data.get("value", []) or []

        db = SessionLocal()
        count = 0
        try:
            for msg in messages:
                sender_obj = (msg.get("from") or {}).get("emailAddress") or {}
                sender_email = (sender_obj.get("address") or "").strip()
                sender_name = (sender_obj.get("name") or sender_email or "Unknown").strip()

                if allowed_senders_set and sender_email.lower() not in allowed_senders_set:
                    continue

                subject = (msg.get("subject") or "").strip()
                body_obj = msg.get("body") or {}
                body_content = body_obj.get("content") or ""
                body_type = (body_obj.get("contentType") or "").lower()

                text_body = _strip_html(body_content) if body_type == "html" else (body_content or "").strip()
                if not text_body:
                    text_body = "(No body content)"

                subject = _truncate(subject, MAX_SUBJECT_LEN)
                text_body = _truncate(text_body, MAX_MESSAGE_TEXT_LEN)
                received_dt = _parse_dt(msg.get("receivedDateTime"))

                row = ScrapedMessage(
                    source="outlook",
                    sender=f"{sender_name} <{sender_email}>" if sender_email else sender_name,
                    sender_role=owner_role,
                    project_id=None,
                    message_datetime=received_dt,
                    subject=subject,
                    message_text=text_body,
                )
                db.add(row)
                count += 1

            db.commit()
        finally:
            db.close()

        return {"ingested": count}

    except requests.RequestException as e:
        raise Exception(f"Network error calling Microsoft Graph: {e}")
    finally:
        session.close()


def ingest_teams_chats(
    access_token: str,
    owner_role="owner",
    max_chats=10,
    max_messages_per_chat=20,
):
    raw = access_token
    access_token = _normalize_token(access_token)

    print("RAW_HAS_SPACES:", any(c.isspace() for c in raw))
    print("RAW_LEN:", len(raw))
    print("CLEAN_LEN:", len(access_token))
    print("TOKEN_START:", access_token[:20])

    # IMPORTANT: validate token really works
    access_token = _validate_ms_access_token(access_token)

    headers = {"Authorization": f"Bearer {access_token}", "Accept": "application/json"}
    session = requests.Session()
    session.trust_env = False

    try:
        chat_url = "https://graph.microsoft.com/v1.0/me/chats"
        r_chats = session.get(chat_url, headers=headers, params={"$top": max_chats}, timeout=30)
        if not r_chats.ok:
            raise Exception(f"Teams Chat Error: {r_chats.text}")

        chats = r_chats.json().get("value", [])
        db = SessionLocal()
        total_ingested = 0

        try:
            for chat in chats:
                chat_id = chat.get("id")
                if not chat_id:
                    continue

                topic = chat.get("topic") or "Direct Chat"

                msg_url = f"https://graph.microsoft.com/v1.0/me/chats/{chat_id}/messages"
                r_msgs = session.get(msg_url, headers=headers, params={"$top": max_messages_per_chat}, timeout=30)

                if not r_msgs.ok:
                    # don't hard-fail all ingestion if one chat errors
                    continue

                for t_msg in r_msgs.json().get("value", []):
                    if t_msg.get("messageType") != "message":
                        continue

                    raw_html = (t_msg.get("body") or {}).get("content") or ""
                    if not raw_html:
                        continue

                    clean_text = _strip_html(raw_html)

                    row = ScrapedMessage(
                        source="teams",
                        sender=(t_msg.get("from") or {}).get("user", {}).get("displayName", "System"),
                        sender_role=owner_role,
                        message_datetime=_parse_dt(t_msg.get("createdDateTime")),
                        subject=f"Teams: {topic}",
                        message_text=_truncate(clean_text, MAX_MESSAGE_TEXT_LEN),
                    )
                    db.add(row)
                    total_ingested += 1

            db.commit()
        finally:
            db.close()

        return {"ingested": total_ingested}
    finally:
        session.close()


def ingest_sharepoint_ppts(
    access_token: str,
    site_id: str,
    drive_id: str,
    folder_path: str,
    owner_role="owner",
):
    headers = {"Authorization": f"Bearer {access_token}"}
    session = requests.Session()
    session.trust_env = False

    path_suffix = f":/{folder_path}:/children" if folder_path and folder_path != "/" else "/root/children"
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}{path_suffix}"

    try:
        r = session.get(url, headers=headers, timeout=30)
        if not r.ok:
            raise Exception(f"SharePoint Drive Error: {r.text}")

        items = r.json().get("value", [])
        db = SessionLocal()
        count = 0

        try:
            for item in items:
                name = (item.get("name") or "").lower()
                if not (item.get("file") and name.endswith(".pptx")):
                    continue

                download_url = item.get("@microsoft.graph.downloadUrl")
                if not download_url:
                    continue

                file_res = session.get(download_url, timeout=60)
                if not file_res.ok:
                    continue

                extracted_text = _extract_text_from_pptx_bytes(file_res.content)

                row = ScrapedMessage(
                    source="ppt",
                    sender="SharePoint",
                    sender_role=owner_role,
                    message_datetime=_parse_dt(item.get("createdDateTime")),
                    subject=_truncate(item.get("name") or "sharepoint.pptx", MAX_SUBJECT_LEN),
                    message_text=_truncate(extracted_text, MAX_MESSAGE_TEXT_LEN),
                )
                db.add(row)
                count += 1

            db.commit()
        finally:
            db.close()

        return {"ingested": count}
    finally:
        session.close()


# =========================
# FastAPI request models + routes (Swagger-friendly)
# =========================
class OutlookIngestReq(BaseModel):
    access_token: Optional[str] = None
    owner_role: str = "owner"
    max_results: int = 50
    allowed_senders: Optional[List[str]] = None


class TeamsIngestReq(BaseModel):
    access_token: Optional[str] = None
    owner_role: str = "owner"
    max_chats: int = 10
    max_messages_per_chat: int = 20


class SharePointPptIngestReq(BaseModel):
    access_token: Optional[str] = None
    owner_role: str = "owner"
    site_id: str
    drive_id: str
    folder_path: str = "/"


@router.post("/ingest/outlook")
def ingest_outlook_endpoint(req: OutlookIngestReq, authorization: Optional[str] = Header(None)):
    try:
        token = _resolve_token(authorization, req.access_token)
        result = ingest_outlook_mail(
            token,
            owner_role=req.owner_role,
            max_results=req.max_results,
            allowed_senders=req.allowed_senders,
        )
        return {"status": "ok", "detail": "Outlook ingestion complete", "result": result}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Outlook ingest failed: {e}")


@router.post("/ingest/teams")
def ingest_teams_endpoint(req: TeamsIngestReq, authorization: Optional[str] = Header(None)):
    try:
        token = _resolve_token(authorization, req.access_token)
        result = ingest_teams_chats(
            token,
            owner_role=req.owner_role,
            max_chats=req.max_chats,
            max_messages_per_chat=req.max_messages_per_chat,
        )
        return {"status": "ok", "detail": "Teams ingestion complete", "result": result}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Teams ingest failed: {e}")


@router.post("/ingest/sharepoint-ppt")
def ingest_sharepoint_ppt_endpoint(req: SharePointPptIngestReq, authorization: Optional[str] = Header(None)):
    try:
        token = _resolve_token(authorization, req.access_token)
        result = ingest_sharepoint_ppts(
            token,
            site_id=req.site_id,
            drive_id=req.drive_id,
            folder_path=req.folder_path,
            owner_role=req.owner_role,
        )
        return {"status": "ok", "detail": "SharePoint PPT ingestion complete", "result": result}
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"SharePoint PPT ingest failed: {e}")