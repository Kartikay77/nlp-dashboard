# backend/app/ingestion/ppt_ingest.py
import io
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from app.db import SessionLocal
from app.models import ScrapedMessage

def _extract_text_from_presentation(prs: Presentation) -> str:
    """Shared logic to pull text from a Presentation object."""
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    chunks.append(f"[Slide {i}] {txt}")
    return "\n".join(chunks).strip()

def _truncate_body(body: str, max_len: int = 20000) -> str:
    if len(body) > max_len:
        return body[:max_len] + "\n\n...[truncated]"
    return body

def extract_ppt_text(ppt_path: str) -> str:
    prs = Presentation(ppt_path)
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    chunks.append(f"[Slide {i}] {txt}")
    return "\n".join(chunks).strip()


def ingest_ppt_folder(folder_path: str, owner_role: str = "owner"):
    db = SessionLocal()
    ingested = 0
    skipped = 0
    failed = 0

    try:
        for p in Path(folder_path).glob("*.pptx"):
            try:
                body = extract_ppt_text(str(p))
                if not body:
                    skipped += 1
                    continue

                # Prevent huge text blobs
                MAX_LEN = 20000
                if len(body) > MAX_LEN:
                    body = body[:MAX_LEN] + "\n\n...[truncated]"

                # MVP dedupe by source + subject (filename-based)
                subject = f"PPT: {p.stem}"
                exists = db.query(ScrapedMessage).filter(
                    ScrapedMessage.source == "ppt",
                    ScrapedMessage.subject == subject
                ).first()

                if exists:
                    skipped += 1
                    continue

                msg = ScrapedMessage(
                    source="ppt",
                    sender="PPT Import",
                    sender_role=owner_role,
                    project_id=None,
                    message_datetime=datetime.fromtimestamp(p.stat().st_mtime),
                    subject=subject,
                    message_text=body,
                )
                db.add(msg)
                ingested += 1

            except Exception:
                failed += 1
                continue

        db.commit()
    finally:
        db.close()

    return {"ingested": ingested, "skipped": skipped, "failed": failed}

def ingest_sharepoint_ppts(access_token: str, site_id: str, drive_id: str, folder_path: str, owner_role="owner"):
    """
    Downloads .pptx files from SharePoint/OneDrive via Microsoft Graph 
    and extracts text for NLP analysis.
    """
    import requests
    headers = {"Authorization": f"Bearer {access_token}"}
    session = requests.Session()
    session.trust_env = False
    
    # Construct Graph API URL for SharePoint children
    path_suffix = f":/{folder_path}:/children" if folder_path and folder_path != "/" else "/root/children"
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}{path_suffix}"
    
    db = SessionLocal()
    count = 0
    
    try:
        r = session.get(url, headers=headers)
        if not r.ok: 
            raise Exception(f"SharePoint API Error: {r.text}")
        
        items = r.json().get("value", [])
        for item in items:
            # Only process PowerPoint files
            if item.get('file') and item['name'].lower().endswith('.pptx'):
                download_url = item.get('@microsoft.graph.downloadUrl')
                if not download_url: continue
                
                # Download file content into memory
                file_res = session.get(download_url)
                if file_res.ok:
                    # Load PPT from bytes
                    ppt_file = io.BytesIO(file_res.content)
                    prs = Presentation(ppt_file)
                    extracted_text = _extract_text_from_presentation(prs)
                    
                    if not extracted_text: continue

                    # Deduplication check
                    subject = f"SharePoint PPT: {item['name']}"
                    exists = db.query(ScrapedMessage).filter(
                        ScrapedMessage.source == "ppt",
                        ScrapedMessage.subject == subject
                    ).first()
                    
                    if not exists:
                        row = ScrapedMessage(
                            source="ppt",
                            sender="SharePoint",
                            sender_role=owner_role,
                            message_datetime=datetime.fromisoformat(item['createdDateTime'].replace("Z", "+00:00")),
                            subject=subject,
                            message_text=_truncate_body(extracted_text),
                        )
                        db.add(row)
                        count += 1
        db.commit()
    finally:
        db.close()
        session.close()
        
    return {"ingested": count}